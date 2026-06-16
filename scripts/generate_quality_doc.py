#!/usr/bin/env python3
"""
质量文档生成工具 - 支持Word/PowerPoint/Excel格式
基于python-docx, python-pptx, openpyxl

Usage:
  python3 generate_quality_doc.py --type docx --template 8d_report --data data.json --output output.docx
  python3 generate_quality_doc.py --type pptx --template 8d_presentation --data data.json --output output.pptx
  python3 generate_quality_doc.py --type xlsx --template control_plan --data data.json --output output.xlsx
"""

import argparse
import json
import os
from datetime import datetime


# ============= Word文档生成 =============
def generate_docx(template, data, output_path):
    """生成Word文档"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
    except ImportError:
        print("错误: python-docx未安装，请运行: pip install python-docx")
        return
    
    doc = Document()
    
    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = '宋体'
    font.size = Pt(10.5)
    
    if template == '8d_report':
        _generate_8d_docx(doc, data)
    elif template == 'vda63_audit_report':
        _generate_vda63_docx(doc, data)
    elif template == 'quality_plan':
        _generate_quality_plan_docx(doc, data)
    elif template == 'work_instruction':
        _generate_work_instruction_docx(doc, data)
    else:
        _generate_generic_docx(doc, template, data)
    
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    doc.save(output_path)
    print(f"Word文档已保存到: {output_path}")


def _generate_8d_docx(doc, data):
    """生成8D报告Word文档"""
    # 标题
    doc.add_heading('8D问题解决报告', level=0)
    
    # 基本信息
    info_table = doc.add_table(rows=6, cols=4)
    info_table.style = 'Table Grid'
    info_data = [
        ['报告编号:', data.get('report_id', ''), '报告日期:', data.get('report_date', datetime.now().strftime('%Y-%m-%d'))],
        ['客户名称:', data.get('customer', ''), '零件编号:', data.get('part_number', '')],
        ['零件名称:', data.get('part_name', ''), '问题发生日期:', data.get('incident_date', '')],
        ['团队领导:', data.get('team_leader', ''), '截止日期:', data.get('deadline', '')],
        ['问题来源:', data.get('source', '客户投诉'), '严重度:', data.get('severity', '')],
        ['当前状态:', data.get('status', '进行中'), '', ''],
    ]
    for i, row_data in enumerate(info_data):
        for j, cell_text in enumerate(row_data):
            cell = info_table.cell(i, j)
            cell.text = cell_text
            if j % 2 == 0:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True
    
    doc.add_paragraph()
    
    # D0-D8 各步骤
    steps = [
        ('D0: 准备工作', 'preparation'),
        ('D1: 成立团队', 'team'),
        ('D2: 描述问题', 'problem_description'),
        ('D3: 临时遏制措施', 'containment'),
        ('D4: 根本原因分析', 'root_cause'),
        ('D5: 永久纠正措施', 'corrective_action'),
        ('D6: 实施纠正措施', 'implementation'),
        ('D7: 预防再发生', 'prevention'),
        ('D8: 团队认可与关闭', 'closure'),
    ]
    
    for step_title, data_key in steps:
        doc.add_heading(step_title, level=1)
        content = data.get(data_key, '待填写...')
        if isinstance(content, list):
            for item in content:
                doc.add_paragraph(str(item), style='List Bullet')
        else:
            doc.add_paragraph(str(content))


def _generate_vda63_docx(doc, data):
    """生成VDA 6.3审核报告"""
    doc.add_heading('VDA 6.3 过程审核报告', level=0)
    
    info_table = doc.add_table(rows=5, cols=4)
    info_table.style = 'Table Grid'
    info_data = [
        ['审核编号:', data.get('audit_id', ''), '审核日期:', data.get('audit_date', '')],
        ['供应商名称:', data.get('supplier', ''), '供应商编号:', data.get('supplier_id', '')],
        ['产品/过程:', data.get('product', ''), '审核员:', data.get('auditor', '')],
        ['审核范围:', data.get('scope', ''), '审核类型:', data.get('audit_type', '初始审核')],
        ['总体评分:', data.get('total_score', ''), '评级:', data.get('rating', '')],
    ]
    for i, row_data in enumerate(info_data):
        for j, cell_text in enumerate(row_data):
            info_table.cell(i, j).text = cell_text
    
    doc.add_paragraph()
    
    # P1-P7各过程要素
    elements = ['P1', 'P2', 'P3', 'P4', 'P5', 'P6', 'P7']
    for element in elements:
        element_data = data.get(element, {})
        doc.add_heading(f'{element}: {element_data.get("name", "")}', level=1)
        doc.add_paragraph(f'评分: {element_data.get("score", "")}%')
        
        questions = element_data.get('questions', [])
        if questions:
            q_table = doc.add_table(rows=len(questions) + 1, cols=4)
            q_table.style = 'Table Grid'
            headers = ['问题编号', '问题描述', '评分', '备注']
            for j, h in enumerate(headers):
                q_table.cell(0, j).text = h
            for i, q in enumerate(questions):
                q_table.cell(i + 1, 0).text = q.get('id', '')
                q_table.cell(i + 1, 1).text = q.get('description', '')
                q_table.cell(i + 1, 2).text = str(q.get('score', ''))
                q_table.cell(i + 1, 3).text = q.get('remark', '')


def _generate_quality_plan_docx(doc, data):
    """生成质量计划"""
    doc.add_heading('产品质量计划', level=0)
    doc.add_paragraph('项目名称: ' + data.get('project_name', ''))
    doc.add_paragraph('客户: ' + data.get('customer', ''))
    doc.add_paragraph('编制日期: ' + data.get('date', datetime.now().strftime('%Y-%m-%d')))
    
    # APQP五阶段
    phases = ['阶段1: 计划和确定项目', '阶段2: 产品设计和开发', 
              '阶段3: 过程设计和开发', '阶段4: 产品和过程确认',
              '阶段5: 反馈、评定和纠正措施']
    
    for phase in phases:
        phase_key = phase.split(':')[0].replace('阶段', 'phase')
        doc.add_heading(phase, level=1)
        phase_data = data.get(phase_key, {})
        doc.add_paragraph(str(phase_data.get('content', '待填写...')))


def _generate_work_instruction_docx(doc, data):
    """生成作业指导书"""
    doc.add_heading('作业指导书', level=0)
    
    info_table = doc.add_table(rows=5, cols=4)
    info_table.style = 'Table Grid'
    info_data = [
        ['文件编号:', data.get('doc_id', ''), '版本:', data.get('version', '1.0')],
        ['工序名称:', data.get('process_name', ''), '工序编号:', data.get('process_id', '')],
        ['适用产品:', data.get('product', ''), '编制日期:', data.get('date', '')],
        ['编制:', data.get('author', ''), '批准:', data.get('approver', '')],
        ['设备/工装:', data.get('equipment', ''), '特殊特性:', data.get('special_chars', '')],
    ]
    for i, row_data in enumerate(info_data):
        for j, cell_text in enumerate(row_data):
            info_table.cell(i, j).text = cell_text
    
    doc.add_paragraph()
    
    # 作业步骤
    steps = data.get('steps', [])
    if steps:
        doc.add_heading('作业步骤', level=1)
        for i, step in enumerate(steps, 1):
            doc.add_heading(f'步骤{i}: {step.get("name", "")}', level=2)
            doc.add_paragraph(f'操作说明: {step.get("instruction", "")}')
            doc.add_paragraph(f'注意事项: {step.get("caution", "")}')
            if step.get('parameters'):
                doc.add_paragraph('过程参数:')
                for param in step['parameters']:
                    doc.add_paragraph(f'  {param.get("name", "")}: {param.get("value", "")} {param.get("unit", "")}', style='List Bullet')


def _generate_generic_docx(doc, template, data):
    """生成通用Word文档"""
    doc.add_heading(template.replace('_', ' ').title(), level=0)
    for key, value in data.items():
        doc.add_heading(key, level=1)
        if isinstance(value, list):
            for item in value:
                doc.add_paragraph(str(item), style='List Bullet')
        elif isinstance(value, dict):
            for k, v in value.items():
                doc.add_paragraph(f'{k}: {v}')
        else:
            doc.add_paragraph(str(value))


# ============= PPT演示文稿生成 =============
def generate_pptx(template, data, output_path):
    """生成PowerPoint演示文稿"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        print("错误: python-pptx未安装，请运行: pip install python-pptx")
        return
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    if template == '8d_presentation':
        _generate_8d_pptx(prs, data)
    elif template == 'quality_monthly':
        _generate_quality_monthly_pptx(prs, data)
    else:
        _generate_generic_pptx(prs, template, data)
    
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    prs.save(output_path)
    print(f"PPT演示文稿已保存到: {output_path}")


def _add_title_slide(prs, title, subtitle=''):
    """添加标题页"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    # 副标题
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(102, 102, 102)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def _add_content_slide(prs, title, bullets):
    """添加内容页"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    # 内容
    contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5.5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f'• {bullet}'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
    return slide


def _generate_8d_pptx(prs, data):
    """生成8D汇报PPT"""
    _add_title_slide(prs, '8D问题解决汇报', data.get('problem_title', ''))
    _add_content_slide(prs, 'D2: 问题描述', [
        data.get('problem_description', '待填写'),
        f"问题来源: {data.get('source', '')}",
        f"影响范围: {data.get('impact', '')}",
    ])
    _add_content_slide(prs, 'D3: 临时遏制措施', [
        data.get('containment', '待填写'),
    ])
    _add_content_slide(prs, 'D4: 根本原因分析', [
        data.get('root_cause', '待填写'),
    ])
    _add_content_slide(prs, 'D5-D6: 纠正措施', [
        data.get('corrective_action', '待填写'),
    ])
    _add_content_slide(prs, 'D7: 预防措施', [
        data.get('prevention', '待填写'),
    ])


def _generate_quality_monthly_pptx(prs, data):
    """生成质量月报PPT"""
    _add_title_slide(prs, '质量月度汇报', data.get('period', ''))
    _add_content_slide(prs, '质量指标概览', [
        f"PPM: {data.get('ppm', 'N/A')}",
        f"FPY: {data.get('fpy', 'N/A')}%",
        f"Cpk: {data.get('cpk', 'N/A')}",
        f"客户投诉: {data.get('complaints', 'N/A')}次",
    ])
    _add_content_slide(prs, '问题与措施', data.get('issues', ['待填写']))


def _generate_generic_pptx(prs, template, data):
    """生成通用PPT"""
    _add_title_slide(prs, template.replace('_', ' ').title())
    for key, value in data.items():
        if isinstance(value, list):
            _add_content_slide(prs, key, value)
        elif isinstance(value, str):
            _add_content_slide(prs, key, [value])


# ============= Excel表格生成 =============
def generate_xlsx(template, data, output_path):
    """生成Excel表格"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError:
        print("错误: openpyxl未安装，请运行: pip install openpyxl")
        return
    
    wb = Workbook()
    ws = wb.active
    
    # 样式
    header_font = Font(name='黑体', size=11, bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='003366', end_color='003366', fill_type='solid')
    center_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
    left_align = Alignment(horizontal='left', vertical='center', wrap_text=True)
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )
    
    if template == 'control_plan':
        _generate_control_plan_xlsx(ws, data, header_font, header_fill, center_align, left_align, thin_border)
    elif template in ('pfmea_sheet', 'dfmea_sheet'):
        _generate_fmea_xlsx(ws, data, header_font, header_fill, center_align, left_align, thin_border)
    else:
        _generate_generic_xlsx(ws, template, data, header_font, header_fill, center_align, left_align, thin_border)
    
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    wb.save(output_path)
    print(f"Excel表格已保存到: {output_path}")


def _generate_control_plan_xlsx(ws, data, header_font, header_fill, center_align, left_align, thin_border):
    """生成控制计划Excel"""
    ws.title = "控制计划"
    
    headers = [
        '过程编号', '过程名称/操作描述', '机器/装置/工装',
        '特性编号', '产品', '过程', '特殊特性分类',
        '产品/过程规格/公差', '评价/测量技术', '样本大小', '频率',
        '控制方法', '反应计划'
    ]
    
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    
    widths = [10, 20, 15, 10, 15, 15, 12, 18, 15, 10, 10, 15, 15]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    
    # 填充数据
    rows = data.get('rows', [])
    for row_idx, row_data in enumerate(rows, 2):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.border = thin_border
            cell.alignment = left_align


def _generate_fmea_xlsx(ws, data, header_font, header_fill, center_align, left_align, thin_border):
    """生成FMEA Excel表格"""
    ws.title = "FMEA"
    
    headers = [
        'FMEA编号', '项目', '过程步骤/设计要素',
        '功能/要求', '失效模式', '失效影响', 'S',
        '失效原因', 'O', '现有预防控制', '现有探测控制', 'D',
        'AP', '建议措施', '责任人/日期', '措施状态',
        '措施后S', '措施后O', '措施后D', '措施后AP'
    ]
    
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    
    for i in range(1, len(headers) + 1):
        ws.column_dimensions[get_column_letter(i)].width = 15


def _generate_generic_xlsx(ws, template, data, header_font, header_fill, center_align, left_align, thin_border):
    """生成通用Excel表格"""
    ws.title = template[:31]
    
    if isinstance(data, dict):
        headers = list(data.keys())
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = center_align
            cell.border = thin_border
        
        values = [str(data.get(h, '')) for h in headers]
        for col, value in enumerate(values, 1):
            cell = ws.cell(row=2, column=col, value=value)
            cell.border = thin_border
            cell.alignment = left_align
    elif isinstance(data, list):
        for row_idx, item in enumerate(data):
            if isinstance(item, dict):
                if row_idx == 0:
                    headers = list(item.keys())
                    for col, header in enumerate(headers, 1):
                        cell = ws.cell(row=1, column=col, value=header)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = center_align
                        cell.border = thin_border
                
                for col, (key, value) in enumerate(item.items(), 1):
                    cell = ws.cell(row=row_idx + 2, column=col, value=str(value))
                    cell.border = thin_border
                    cell.alignment = left_align


def main():
    parser = argparse.ArgumentParser(description='质量文档生成工具')
    parser.add_argument('--type', required=True, choices=['docx', 'pptx', 'xlsx'], help='文档类型')
    parser.add_argument('--template', required=True, help='文档模板名')
    parser.add_argument('--data', required=True, help='数据文件路径(JSON)')
    parser.add_argument('--output', required=True, help='输出文件路径')
    
    args = parser.parse_args()
    
    with open(args.data, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    if args.type == 'docx':
        generate_docx(args.template, data, args.output)
    elif args.type == 'pptx':
        generate_pptx(args.template, data, args.output)
    elif args.type == 'xlsx':
        generate_xlsx(args.template, data, args.output)


if __name__ == '__main__':
    main()
